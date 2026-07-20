#!/usr/bin/env python3
"""flowcast PPT export (B-out) — 흐름도 뷰 JSON → 편집 가능한 네이티브 .pptx.

python-pptx 는 **export 경로 전용 선택적 의존성**이다(코어 render/import 는 stdlib 유지).
미설치 시 친절히 안내하고 종료한다.

좌표는 render.py 의 배치 로직(`_c_rect`·`_t_rect`·`_edge_pt`·`layout_sequence`)을 **그대로
재사용**해 HTML/PDF 와 동일한 상대 위치로 도형을 찍는다(재구현 금지). px → EMU(×9525) 변환.

슬라이드 캔버스(기본 `wide` = 1920×1080px, 16:9): 콘텐츠를 uniform scale 로 fit(업스케일
포함)하고 중앙 배치한다. 폰트도 같은 배율로 스케일(0.5pt 반올림, 최소 6pt).
`--slide-size auto` 는 기존 content-fit(콘텐츠 크기 = 슬라이드 크기) 동작.

사용법:
    python3 scripts/pptx_export.py {view.json} [-o {out.pptx}] [--slide-size wide|auto|{W}x{H}]

지원: sequence · component · topology 3뷰 (view 필드로 디스패치, 미지정=sequence).
- component: 노드 → 둥근 사각형 + 텍스트, 존 → 배경 사각형, 엣지 → 커넥터 + 라벨
- topology: 위와 동일 + 세그먼트 라벨은 번호 배지(원) + 하단 "흐름 설명" 범례 (render.py 패리티)
- sequence: 액터 박스 + 라이프라인 + 액티베이션 바 + 메시지 커넥터(kind별 색) + note 박스
"""

import argparse
import importlib.util
import json
import sys
from pathlib import Path

PX_TO_EMU = 9525   # 1px @96dpi
PAD_PX = 20        # auto(content-fit) 슬라이드 여백(px)
FIT_MARGIN = 60    # 고정 캔버스 fit 여백(px)
# 상단/하단 크롬 — 슬라이드 절대 px(콘텐츠 배율과 무관)로 고정해 뷰 간 상단 일관성 확보.
# render.py .hero(항상 flush 좌상단) 대응: pptx 는 fit-centering 아티팩트로 어긋나므로 고정 앵커로 정규화.
HEADER_LEFT = 60   # 제목 좌측 고정 (= FIT_MARGIN, 콘텐츠 영역 좌측과 정렬)
HEADER_TOP = 42    # 제목 상단 고정
HEADER_BAND = 150  # 콘텐츠 상단 경계(px) — 제목 밴드/본문 분리선
FOOTER_BAND = 46   # 하단 푸터 밴드 높이(px) — 본문/푸터 분리선 = SH - FOOTER_BAND
HEADER_RULE_Y = HEADER_BAND - 6   # 제목/본문 구분 실선 y (콘텐츠 살짝 위)
TITLE_PT = {"eyebrow": 12.5, "h1": 22, "sub": 12}   # 제목 고정 폰트(pt) — 뷰 무관 통일
RULE_RGB = (0x2C, 0x7A, 0x7B)     # 구분선 색 (flowcast comp line teal)
RULE_PX = 1.5                     # 구분선 두께(px)
SEQ_PAGE_MIN_SCALE = 0.8   # 긴 sequence: 한 장 scale 이 이 값 미만이면 페이지 분할(가독 하한)
SLIDE_PRESETS = {"wide": (1920, 1080)}   # 캔버스 프리셋 (px) — 기본 wide
ACCENT = (0x1F, 0x6F, 0xD0)              # --accent (topology 번호 배지)

FILL = {"comp": (0xE6, 0xF4, 0xF1), "ext": (0xFD, 0xF3, 0xE7)}
LINE = {"comp": (0x2C, 0x7A, 0x7B), "ext": (0xB7, 0x79, 0x1F)}

# topology kind: srv(기본) · ext(외부, 앰버) · gear(장비, 점선) · fw(방화벽, --warn 굵은 테두리)
# · l4(VIP/LB, --accent 점선). 키 집합은 render.py TOPO_KINDS 와 일치해야 한다(기동 시 대조).
TOPO_FILL = {"srv": (0xF8, 0xFA, 0xFC), "ext": (0xFD, 0xF3, 0xE7), "gear": (0xF8, 0xFA, 0xFC),
             "fw": (0xF8, 0xFA, 0xFC), "l4": (0xF8, 0xFA, 0xFC)}
TOPO_LINE = {"srv": (0x64, 0x74, 0x8B), "ext": (0xB7, 0x79, 0x1F), "gear": (0x94, 0xA3, 0xB8),
             "fw": (0x8A, 0x62, 0x10), "l4": (0x1F, 0x6F, 0xD0)}
TOPO_DASH = {"gear", "l4"}        # HTML 의 stroke-dasharray 대응
TOPO_LINE_W = {"fw": 1.6}         # 그 외 기본(1.0) — HTML .topo-fw{stroke-width:1.6}

# sequence kind: req/self(accent) · res/relay(muted) — render.py LIGHT 테마 흰배경 솔리드 근사
SEQ_ARROW = {"req": (0x1F, 0x6F, 0xD0), "self": (0x1F, 0x6F, 0xD0),
             "res": (0x54, 0x66, 0x7E), "relay": (0x54, 0x66, 0x7E)}
SEQ_ACTOR_FILL = (0xEF, 0xF5, 0xFC)   # --zone-bg over white
SEQ_ACTOR_LINE = (0x1F, 0x6F, 0xD0)   # --accent
SEQ_ACT_FILL = (0xE4, 0xEE, 0xF9)     # --act-bg over white
SEQ_ACT_LINE = (0xB7, 0xD1, 0xF0)     # --act-bd over white
SEQ_NOTE_FILL = (0xFA, 0xF6, 0xEE)    # --note-bg over white
SEQ_NOTE_LINE = (0xE0, 0xC9, 0x98)    # --note-bd over white
SEQ_TEXT = (0x1B, 0x26, 0x35)         # --text
SEQ_MUTED = (0x54, 0x66, 0x7E)        # --muted
SEQ_WARN = (0x8A, 0x62, 0x10)         # --warn
SEQ_LIFELINE = (0x9C, 0xA3, 0xAF)     # 회색 점선


def _load_render():
    spec = importlib.util.spec_from_file_location(
        "flowcast_render", Path(__file__).parent / "render.py")
    m = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(m)
    return m


def _import_pptx():
    try:
        import pptx  # noqa: F401
        return True
    except ImportError:
        return False


def _parse_slide_size(opt):
    """'wide'|'auto'|'{W}x{H}' → (W, H) px. auto 는 None(content-fit)."""
    if opt == "auto":
        return None
    if opt in SLIDE_PRESETS:
        return SLIDE_PRESETS[opt]
    try:
        w, h = str(opt).lower().split("x", 1)
        return int(w), int(h)
    except ValueError:
        raise SystemExit(f"--slide-size 형식 오류: {opt!r} (wide|auto|{{W}}x{{H}})")


def _fit(content_w, content_h, size):
    """콘텐츠(px)를 캔버스에 uniform scale 로 fit — (slide_w, slide_h, scale, dx, dy).

    콘텐츠는 상단 HEADER_BAND · 하단 FOOTER_BAND(절대 px) 사이 영역에 fit·중앙배치된다.
    dx/dy 는 콘텐츠-px 단위 오프셋: 스케일된 emu 변환에서 (좌표 + dx) * scale = 슬라이드-px.
    호출부는 raw content_w/content_h 를 넘긴다(타이틀 밴드 예약은 여기서 처리).
    size=None(auto) 은 content-fit — scale 1, 콘텐츠 = 캔버스, 헤더/푸터 밴드만 상하로 예약.
    """
    cw, ch = max(content_w, 1), max(content_h, 1)
    if size is None:
        W = cw + 2 * PAD_PX
        H = HEADER_BAND + ch + FOOTER_BAND
        return W, H, 1.0, (W - cw) / 2, HEADER_BAND     # 콘텐츠 좌우 중앙 · 헤더 밴드 아래
    W, H = size
    avail_w = W - 2 * FIT_MARGIN
    avail_h = H - HEADER_BAND - FOOTER_BAND
    s = min(avail_w / cw, avail_h / ch)
    dx = (FIT_MARGIN + (avail_w - cw * s) / 2) / s      # [FIT_MARGIN, W-FIT_MARGIN] 중앙
    dy = (HEADER_BAND + (avail_h - ch * s) / 2) / s     # [HEADER_BAND, H-FOOTER_BAND] 중앙
    return W, H, s, dx, dy


def _fpt(base, s):
    """폰트 pt 를 배율 s 로 스케일 — 0.5pt 반올림, 최소 6pt."""
    return max(6.0, round(base * s * 2) / 2)


def _fill_lines(tf, name, s, size=10, sub_size=8,
                color=(0x1A, 0x20, 0x2C), sub_color=(0x4B, 0x55, 0x63), bold=True):
    """멀티라인 이름을 문단별로 채움 — 모든 문단에 명시적 스타일(첫 줄 base·bold, 이후 sub).

    `tf.text = "a\\nb"` 는 문단을 쪼개는데 기존 코드가 paragraphs[0]만 스타일링해
    둘째 줄부터 템플릿 기본값(18pt)으로 새던 버그의 공통 해소 지점.
    """
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    for i, ln in enumerate(str(name).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = ln
        first = i == 0
        r.font.size = Pt(_fpt(size if first else sub_size, s))
        r.font.bold = bold and first
        r.font.color.rgb = RGBColor(*(color if first else sub_color))


def _rule(shapes, x1, x2, y, rgb=RULE_RGB, thick=RULE_PX):
    """가로 구분 실선 — 얇은 사각형(AUTO_SHAPE)으로 그린다.

    add_connector(LINE) 를 쓰면 커넥터 수 테스트(=엣지/라이프라인 수)와 충돌하므로 금지.
    좌표는 슬라이드 절대 px(× PX_TO_EMU, 콘텐츠 배율 s 무관).
    """
    from pptx.util import Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    a = lambda px: Emu(int(round(px * PX_TO_EMU)))
    bar = shapes.add_shape(MSO_SHAPE.RECTANGLE, a(x1), a(y - thick / 2), a(x2 - x1), a(thick))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(*rgb)
    bar.line.fill.background()
    bar.shadow.inherit = False
    return bar


def _draw_chrome(shapes, SW, SH, system, title, source, idx, total):
    """상단 제목 + 제목/본문·본문/푸터 구분 실선 + 페이지 — 모두 슬라이드 절대 px 고정.

    render.py .hero(flush 좌상단) 대응. 제목은 고정 앵커(HEADER_LEFT/TOP)+고정 폰트로 그려
    콘텐츠 배율 s 와 무관하게 뷰·슬라이드 상단이 일관된다. 구분선은 _rule(사각형).
    """
    from pptx.util import Emu, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    a = lambda px: Emu(int(round(px * PX_TO_EMU)))

    # 제목 밴드 (eyebrow · h1 · sub) — 고정 좌상단·고정 폰트
    tb = shapes.add_textbox(a(HEADER_LEFT), a(HEADER_TOP),
                            a(SW - 2 * HEADER_LEFT), a(HEADER_BAND - HEADER_TOP))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = tf.margin_top = 0
    specs = []
    if system:
        specs.append((f"{system} · IF 흐름도", TITLE_PT["eyebrow"], (0x1F, 0x6F, 0xD0), True))
    specs.append((str(title), TITLE_PT["h1"], (0x1B, 0x26, 0x35), True))
    if source:
        specs.append((str(source), TITLE_PT["sub"], (0x54, 0x66, 0x7E), False))
    for i, (text, size, rgb, bold) in enumerate(specs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = RGBColor(*rgb)

    # 제목/본문 · 본문/푸터 구분 실선
    _rule(shapes, FIT_MARGIN, SW - FIT_MARGIN, HEADER_RULE_Y)
    _rule(shapes, FIT_MARGIN, SW - FIT_MARGIN, SH - FOOTER_BAND)

    # 푸터: system(좌) · 페이지(우) — 회색, 작게
    fy = SH - FOOTER_BAND + 6
    if system:
        lt = shapes.add_textbox(a(HEADER_LEFT), a(fy), a(SW - 2 * HEADER_LEFT - 160), a(28))
        ltf = lt.text_frame
        ltf.margin_left = ltf.margin_top = 0
        lr = ltf.paragraphs[0].add_run()
        lr.text = str(system)
        lr.font.size = Pt(11)
        lr.font.color.rgb = RGBColor(0x54, 0x66, 0x7E)
    pg = shapes.add_textbox(a(SW - FIT_MARGIN - 160), a(fy), a(160), a(28))
    ptf = pg.text_frame
    ptf.margin_left = ptf.margin_right = ptf.margin_top = 0
    pp = ptf.paragraphs[0]
    pp.alignment = PP_ALIGN.RIGHT
    pr = pp.add_run()
    pr.text = f"{idx} / {total}"
    pr.font.size = Pt(11)
    pr.font.color.rgb = RGBColor(0x54, 0x66, 0x7E)


def _scene_geometry(R, scenario):
    """시나리오의 노드 사각형·존 박스·전체 bounding box (px)."""
    nodes = scenario.get("nodes") or []
    zones = scenario.get("zones") or []
    rects = {n["id"]: R._c_rect(n) for n in nodes}
    xs, ys = [], []
    for x, y, w, h in rects.values():
        xs += [x, x + w]
        ys += [y, y + h]
    zone_boxes = []
    for z in zones:
        members = [n["id"] for n in nodes if n.get("zone") == z["id"]]
        mr = [rects[m] for m in members]
        if not mr:
            continue
        zx1 = min(r[0] for r in mr) - R.C_ZONE_PAD
        zy1 = min(r[1] for r in mr) - R.C_ZONE_PAD - R.C_ZONE_LBL
        zx2 = max(r[0] + r[2] for r in mr) + R.C_ZONE_PAD
        zy2 = max(r[1] + r[3] for r in mr) + R.C_ZONE_PAD
        zone_boxes.append((z, zx1, zy1, zx2, zy2))
        xs += [zx1, zx2]
        ys += [zy1, zy2]
    bbox = (min(xs), min(ys), max(xs), max(ys)) if xs else (0, 0, 0, 0)
    return rects, zone_boxes, bbox


def export_component(data, out_path, slide_size="wide"):
    """component 뷰 JSON → .pptx (시나리오 1개 = 슬라이드 1장)."""
    from pptx import Presentation
    from pptx.util import Emu, Pt
    from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
    from pptx.oxml.ns import qn

    R = _load_render()
    scenarios = data.get("scenarios") or []

    # 캔버스 = 시나리오 중 최대 bounding box 를 fit (좌표·크기는 스케일된 emu 로 변환)
    geoms = [_scene_geometry(R, sc) for sc in scenarios]
    max_w = max(((b[2] - b[0]) for _, _, b in geoms), default=400)
    max_h = max(((b[3] - b[1]) for _, _, b in geoms), default=300)
    SW, SH, s, dx, dy = _fit(max_w, max_h, _parse_slide_size(slide_size))
    emu = lambda px: Emu(int(round(px * s * PX_TO_EMU)))

    prs = Presentation()
    prs.slide_width = Emu(SW * PX_TO_EMU)
    prs.slide_height = Emu(SH * PX_TO_EMU)
    blank = prs.slide_layouts[6]

    def _arrow(connector, tail=True, head=False):
        ln = connector.line._get_or_add_ln()
        if head:
            ln.append(ln.makeelement(qn("a:headEnd"), {"type": "triangle"}))
        if tail:
            ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle"}))

    for idx, (scenario, (rects, zone_boxes, bbox)) in enumerate(zip(scenarios, geoms), 1):
        minx, miny = bbox[0], bbox[1]
        ox, oy = -minx + dx, -miny + dy   # 원점 이동 + 중앙 배치(헤더 밴드 아래)
        slide = prs.slides.add_slide(blank)
        shapes = slide.shapes

        _draw_chrome(shapes, SW, SH, data.get("system", ""), scenario.get("title", ""),
                     data.get("source", ""), idx, len(scenarios))

        # 존 배경 먼저(뒤에 깔림)
        for z, zx1, zy1, zx2, zy2 in zone_boxes:
            zb = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  emu(zx1 + ox), emu(zy1 + oy),
                                  emu(zx2 - zx1), emu(zy2 - zy1))
            zb.fill.background()
            zb.line.color.rgb = RGBColor(0x9C, 0xA3, 0xAF)
            zb.line.dash_style = None
            zb.shadow.inherit = False
            tf = zb.text_frame
            tf.text = z["name"]
            tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
            tf.paragraphs[0].runs[0].font.size = Pt(_fpt(9, s))
            tf.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x6B, 0x72, 0x80)
            tf.word_wrap = True
            zb.text_frame.vertical_anchor = MSO_ANCHOR.TOP

        # 엣지 커넥터 — 노드보다 먼저 (render.py z-순서 패리티: 노드가 선을 가림, #25)
        pending_labels = []
        for e in scenario.get("edges") or []:
            r1, r2 = rects.get(e["from"]), rects.get(e["to"])
            if not r1 or not r2:
                continue
            c1 = (r1[0] + r1[2] / 2, r1[1] + r1[3] / 2)
            c2 = (r2[0] + r2[2] / 2, r2[1] + r2[3] / 2)
            ax, ay = R._edge_pt(r1, *c2)
            bx, by = R._edge_pt(r2, *c1)
            conn = shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                        emu(ax + ox), emu(ay + oy),
                                        emu(bx + ox), emu(by + oy))
            conn.line.color.rgb = RGBColor(0x4B, 0x55, 0x63)
            _arrow(conn, tail=True, head=bool(e.get("bidir")))

            parts = []
            if e.get("n") is not None:
                parts.append(f"({e['n']})")
            if e.get("label"):
                parts.append(e["label"])
            proto = f"( {e['protocol']} )" if e.get("protocol") else ""
            if parts or proto:
                mx = (e["lx"] if e.get("lx") is not None else (ax + bx) / 2)
                my = (e["ly"] if e.get("ly") is not None else (ay + by) / 2)
                pending_labels.append((mx, my, " ".join(parts), proto))

        # 노드 (커넥터 위)
        node_by_id = {n["id"]: n for n in scenario.get("nodes") or []}
        for nid, (x, y, w, h) in rects.items():
            nd = node_by_id[nid]
            kind = nd.get("kind", "comp")
            box = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   emu(x + ox), emu(y + oy), emu(w), emu(h))
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(*FILL.get(kind, FILL["comp"]))
            box.line.color.rgb = RGBColor(*LINE.get(kind, LINE["comp"]))
            box.shadow.inherit = False
            tf = box.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            _fill_lines(tf, nd["name"], s, size=11)   # 밝은 fill 위 가독성
            if nd.get("port"):
                pp = tf.add_paragraph()
                pp.alignment = PP_ALIGN.CENTER
                r = pp.add_run()
                r.text = f"Port: {nd['port']}"
                r.font.size = Pt(_fpt(8, s))
                r.font.color.rgb = RGBColor(0x4B, 0x55, 0x63)

        # 엣지 라벨 (번호 인라인 + 프로토콜) — 최상위
        for mx, my, text, proto in pending_labels:
            tb = shapes.add_textbox(emu(mx + ox - 60), emu(my + oy - 12),
                                    emu(140), emu(30))
            tf = tb.text_frame
            tf.word_wrap = True
            tf.text = text
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            tf.paragraphs[0].runs[0].font.size = Pt(_fpt(9, s))
            if proto:
                pp = tf.add_paragraph()
                pp.alignment = PP_ALIGN.CENTER
                r = pp.add_run()
                r.text = proto
                r.font.size = Pt(_fpt(8, s))

    prs.save(str(out_path))
    return len(scenarios)


def export_topology(data, out_path, slide_size="wide"):
    """topology 뷰 JSON → .pptx. nodes/links/zones 공유(모든 슬라이드) + 시나리오별 segments 오버레이.

    세그먼트 라벨은 render.py 패리티 — 엣지엔 원형 번호 배지만, 전문은 하단 "흐름 설명" 범례.
    """
    from pptx import Presentation
    from pptx.util import Emu, Pt
    from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
    from pptx.enum.dml import MSO_LINE_DASH_STYLE
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
    from pptx.oxml.ns import qn

    R = _load_render()
    # kind 팔레트는 render.py TOPO_KINDS 가 단일 진실 — 새 kind 가 추가되면 여기서 드러난다.
    _missing = R.TOPO_KINDS - set(TOPO_FILL)
    if _missing:
        print(f"경고: pptx 팔레트에 없는 topology kind {sorted(_missing)} — srv 로 그려집니다.",
              file=sys.stderr)
    nodes = data.get("nodes") or []
    zones = data.get("zones") or []
    links = data.get("links") or []
    scenarios = data.get("scenarios") or [{"title": data.get("system", "topology")}]
    rects = {n["id"]: R._t_rect(n) for n in nodes}
    node_by_id = {n["id"]: n for n in nodes}

    xs, ys = [], []
    for x, y, w, h in rects.values():
        xs += [x, x + w]
        ys += [y, y + h]
    zone_boxes = []
    for z in zones:
        members = [n["id"] for n in nodes if n.get("zone") == z["id"]]
        mr = [rects[m] for m in members]
        if not mr:
            continue
        zx1 = min(r[0] for r in mr) - R.T_ZONE_PAD
        zy1 = min(r[1] for r in mr) - R.T_ZONE_PAD - R.T_ZONE_LBL
        zx2 = max(r[0] + r[2] for r in mr) + R.T_ZONE_PAD
        zy2 = max(r[1] + r[3] for r in mr) + R.T_ZONE_PAD
        zone_boxes.append((z, zx1, zy1, zx2, zy2))
        xs += [zx1, zx2]
        ys += [zy1, zy2]
    minx, miny = (min(xs), min(ys)) if xs else (0, 0)
    maxx, maxy = (max(xs), max(ys)) if xs else (0, 0)

    # 흐름 설명 범례(시나리오별) — render.py 공용 표 모델로 네이티브 add_table 렌더
    leg_x, leg_y = minx, maxy + 30
    diagram_right = maxx + R.T_BOX_W                     # 노드 오른쪽 끝 근사 (render.py 와 동일 기준)
    legends, leg_max_x, leg_bottom = [], leg_x, leg_y
    for sc in scenarios:
        labelled = [sg for sg in (sc.get("segments") or []) if sg.get("label")]
        if labelled:
            tables, ly_end, lmx = R._topo_legend_tables(labelled, leg_x, leg_y, diagram_right)
        else:
            tables, ly_end, lmx = [], leg_y, leg_x
        legends.append(tables)
        leg_max_x = max(leg_max_x, lmx)
        leg_bottom = max(leg_bottom, ly_end)
    if any(legends):
        maxy = leg_bottom + 8
        maxx = max(maxx, leg_max_x)

    SW, SH, s, dx, dy = _fit(maxx - minx, maxy - miny, _parse_slide_size(slide_size))
    ox, oy = -minx + dx, -miny + dy   # 헤더 밴드 아래 중앙 배치
    emu = lambda px: Emu(int(round(px * s * PX_TO_EMU)))
    prs = Presentation()
    prs.slide_width = Emu(SW * PX_TO_EMU)
    prs.slide_height = Emu(SH * PX_TO_EMU)
    blank = prs.slide_layouts[6]

    def _anchor(r1, r2):
        c1 = (r1[0] + r1[2] / 2, r1[1] + r1[3] / 2)
        c2 = (r2[0] + r2[2] / 2, r2[1] + r2[3] / 2)
        return R._edge_pt(r1, *c2), R._edge_pt(r2, *c1)

    def _badge(bx, by, n, rad=11):
        """원형 번호 배지 (render.py 패리티) — accent 채움 + 흰 숫자. rad=9 는 흐름 설명 표용."""
        b = shapes.add_shape(MSO_SHAPE.OVAL, emu(bx + ox - rad), emu(by + oy - rad),
                             emu(2 * rad), emu(2 * rad))
        b.name = "legbadge" if rad < 11 else "badge"   # 다이어그램 배지 vs 흐름 설명 표 배지 구분
        b.fill.solid()
        b.fill.fore_color.rgb = RGBColor(*ACCENT)
        b.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        b.shadow.inherit = False
        tf = b.text_frame
        tf.word_wrap = False
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = str(n)
        r.font.size = Pt(_fpt(8 if rad >= 11 else 7, s))
        r.font.bold = True
        r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    def _leg_table_style(gf):
        """네이티브 표를 'Table Grid'(얇은 회색 괘선·밴딩 없음)로 — 편집가능 깔끔한 표."""
        tbl = gf._element.graphic.graphicData.tbl
        tblPr = tbl.find(qn("a:tblPr"))
        if tblPr is None:
            tblPr = tbl.makeelement(qn("a:tblPr"), {})
            tbl.insert(0, tblPr)
        tblPr.set("firstRow", "1")
        tblPr.set("bandRow", "0")
        sid = tblPr.find(qn("a:tableStyleId"))
        if sid is None:
            sid = tblPr.makeelement(qn("a:tableStyleId"), {})
            tblPr.append(sid)
        sid.text = "{5940675A-B579-460E-94D1-54222C63F5DA}"   # Table Grid

    def _leg_cell(cell, paras, align=PP_ALIGN.LEFT, fill=None, valign=MSO_ANCHOR.TOP):
        """표 셀 채우기. paras = [(text, size, bold, (r,g,b)), ...] 문단 목록."""
        cell.vertical_anchor = valign
        cell.margin_left = cell.margin_right = emu(4)
        cell.margin_top = cell.margin_bottom = emu(2)
        if fill is None:
            cell.fill.background()
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(*fill)
        tf = cell.text_frame
        tf.word_wrap = True
        for i, (text, size, bold, color) in enumerate(paras):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            r = p.add_run()
            r.text = text
            r.font.size = Pt(_fpt(size, s))
            r.font.bold = bold
            r.font.color.rgb = RGBColor(*color)

    for idx, (scenario, leg_tables) in enumerate(zip(scenarios, legends), 1):
        segments = scenario.get("segments") or []
        slide = prs.slides.add_slide(blank)
        shapes = slide.shapes

        _draw_chrome(shapes, SW, SH, data.get("system", ""), scenario.get("title", ""),
                     data.get("source", ""), idx, len(scenarios))

        for z, zx1, zy1, zx2, zy2 in zone_boxes:
            zb = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, emu(zx1 + ox), emu(zy1 + oy),
                                  emu(zx2 - zx1), emu(zy2 - zy1))
            zb.fill.background()
            zb.line.color.rgb = RGBColor(0x9C, 0xA3, 0xAF)
            zb.shadow.inherit = False
            tf = zb.text_frame
            tf.text = z["name"]
            tf.vertical_anchor = MSO_ANCHOR.TOP
            p0 = tf.paragraphs[0]
            p0.alignment = PP_ALIGN.LEFT   # render.py 존 라벨 top-left 패리티
            p0.runs[0].font.size = Pt(_fpt(9, s))
            p0.runs[0].font.color.rgb = RGBColor(0x6B, 0x72, 0x80)

        # 정적 배선(links) — 화살촉 없는 회색선. 노드보다 먼저 (HTML z-순서 패리티, #25)
        dual_ids = {n["id"] for n in nodes if n.get("dual")}
        for lk in links:
            r1, r2 = rects.get(lk.get("from")), rects.get(lk.get("to"))
            if not r1 or not r2:
                continue
            if lk.get("to") in dual_ids:   # 이중화 서버 진입 — 2박스 양쪽으로 분기
                g = 7
                bh = (r2[3] - g) / 2
                r2_targets = [(r2[0], r2[1], r2[2], bh), (r2[0], r2[1] + bh + g, r2[2], bh)]
            else:
                r2_targets = [r2]
            for rb in r2_targets:
                (ax, ay), (bx, by) = _anchor(r1, rb)
                ln = shapes.add_connector(MSO_CONNECTOR.STRAIGHT, emu(ax + ox), emu(ay + oy),
                                          emu(bx + ox), emu(by + oy))
                ln.line.color.rgb = RGBColor(0x94, 0xA3, 0xB8)

        # 구간 오버레이(segments) — 화살촉 커넥터. 배지는 노드 위에 spread 후 일괄 배치.
        badge_sgs = []
        for sg in segments:
            r1 = rects.get(sg.get("from"))
            if not r1:
                continue
            if sg.get("n") is not None:
                badge_sgs.append(sg)
            to = sg.get("to")
            if sg.get("self") or not to or to not in rects:
                continue   # self·대상없음 커넥터 라우팅은 범위 밖 (배지는 노드 위)
            r2 = rects[to]
            (ax, ay), (bx, by) = _anchor(r1, r2)
            conn = shapes.add_connector(MSO_CONNECTOR.STRAIGHT, emu(ax + ox), emu(ay + oy),
                                        emu(bx + ox), emu(by + oy))
            conn.line.color.rgb = RGBColor(0x33, 0x41, 0x55)
            ln = conn.line._get_or_add_ln()
            ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle"}))

        # 노드 — 커넥터 위 (관통 선을 가림)
        for nid, (x, y, w, h) in rects.items():
            nd = node_by_id[nid]
            kind = nd.get("kind", "srv")
            if kind not in TOPO_FILL:
                print(f"경고: 알 수 없는 topology kind {kind!r} — srv 로 그립니다.", file=sys.stderr)
            fillc = RGBColor(*TOPO_FILL.get(kind, TOPO_FILL["srv"]))
            linec = RGBColor(*TOPO_LINE.get(kind, TOPO_LINE["srv"]))
            dashed = kind in TOPO_DASH
            linew = TOPO_LINE_W.get(kind)
            if nd.get("dual"):   # 이중화(2대) — 그룹 테두리 + 위/아래 분리 2박스 (IP 하나씩)
                g = 7
                bh = (h - g) / 2
                box_lines = R._split_dual_ip(nd["name"])
                grp = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, emu(x + ox - 5), emu(y + oy - 5), emu(w + 10), emu(h + 10))
                grp.fill.background()
                grp.line.color.rgb = RGBColor(0x9C, 0xA3, 0xAF)
                grp.shadow.inherit = False
                for bi, by0 in enumerate((y, y + bh + g)):
                    bx = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, emu(x + ox), emu(by0 + oy), emu(w), emu(bh))
                    bx.fill.solid()
                    bx.fill.fore_color.rgb = fillc
                    bx.line.color.rgb = linec
                    bx.shadow.inherit = False
                    tf = bx.text_frame
                    tf.word_wrap = True
                    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                    _fill_lines(tf, "\n".join(box_lines[bi]), s, size=9)
                continue
            box = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, emu(x + ox), emu(y + oy), emu(w), emu(h))
            box.fill.solid()
            box.fill.fore_color.rgb = fillc
            box.line.color.rgb = linec
            if dashed:
                box.line.dash_style = MSO_LINE_DASH_STYLE.DASH
            if linew:
                box.line.width = Pt(linew)
            box.shadow.inherit = False
            tf = box.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            _fill_lines(tf, nd["name"], s, size=10)

        # 번호 배지 — render.py 공용 헬퍼(_t_badge_geom/_t_spread_badges)로 겹침 자동 회피
        spread = R._t_spread_badges([R._t_badge_geom(sg, rects) for sg in badge_sgs])
        for sg, (bx, by) in zip(badge_sgs, spread):
            _badge(bx, by, sg["n"])

        # 흐름 설명 범례 — 네이티브 add_table (편집가능·정렬된 표). 다단이면 열별 표 N개.
        if leg_tables:
            hb = shapes.add_textbox(emu(leg_x + ox), emu(leg_y + oy - R.T_LEG_LH),
                                    emu(200), emu(R.T_LEG_LH + 6))
            htf = hb.text_frame
            htf.word_wrap = False
            htf.margin_left = htf.margin_right = htf.margin_top = htf.margin_bottom = 0
            hr = htf.paragraphs[0].add_run()
            hr.text = "흐름 설명"
            hr.font.size = Pt(_fpt(9, s))
            hr.font.bold = True
            hr.font.color.rgb = RGBColor(0x54, 0x66, 0x7E)
            HDR_FILL = (0xEC, 0xEF, 0xF3)
            STEP_C, DESC_C, META_C = (0x1B, 0x26, 0x35), (0x1B, 0x26, 0x35), (0x54, 0x66, 0x7E)
            for tb in leg_tables:
                gf = shapes.add_table(1 + len(tb["rows"]), 3,
                                      emu(tb["x"] + ox), emu(tb["y"] + oy),
                                      emu(tb["w"]), emu(tb["total_h"]))
                _leg_table_style(gf)
                table = gf.table
                bw, sw, dw = tb["col_w"]
                table.columns[0].width = emu(bw)
                table.columns[1].width = emu(sw)
                table.columns[2].width = emu(dw)
                _leg_cell(table.cell(0, 0), [("#", 8, True, META_C)], PP_ALIGN.CENTER, HDR_FILL, MSO_ANCHOR.MIDDLE)
                _leg_cell(table.cell(0, 1), [("단계", 8, True, META_C)], PP_ALIGN.CENTER, HDR_FILL, MSO_ANCHOR.MIDDLE)
                _leg_cell(table.cell(0, 2), [("설명 · 기술", 8, True, META_C)], PP_ALIGN.LEFT, HDR_FILL, MSO_ANCHOR.MIDDLE)
                table.rows[0].height = emu(tb["header_h"])
                for ri, row in enumerate(tb["rows"], 1):
                    num = "" if row["n"] is None else str(row["n"])
                    _leg_cell(table.cell(ri, 0), [(num, 9, True, ACCENT)], PP_ALIGN.CENTER, None, MSO_ANCHOR.MIDDLE)
                    _leg_cell(table.cell(ri, 1), [(row["step"] or "", 8.5, True, STEP_C)],
                              PP_ALIGN.CENTER, None, MSO_ANCHOR.MIDDLE)
                    paras = [(ln, 8.5, False, DESC_C) for ln in row["desc"]] + \
                            [(ln, 7.5, False, META_C) for ln in row["meta"]]
                    _leg_cell(table.cell(ri, 2), paras or [("", 8.5, False, DESC_C)])
                    table.rows[ri].height = emu(row["h"])

    prs.save(str(out_path))
    return len(scenarios)


def _paginate_sequences(R, data, scenarios, size):
    """긴 시나리오를 읽히는 크기로 페이지 분할 — 한 장 scale < MIN 이면 스텝을 나눠 여러 슬라이드로.

    분할 규칙:
      - 트리거: 시나리오 전체를 한 장에 넣을 때 s < SEQ_PAGE_MIN_SCALE (짧으면 그대로 1장).
      - 예산: avail_h / MIN_SCALE 높이까지 스텝을 그리디로 채우고 **스텝 경계에서만** 넘김
        (note·self 중간 절단 방지). 액터/존/라이프라인은 export 가 슬라이드마다 재렌더 → 연장선.
      - step.n 은 원본 유지(번호 연속). 제목에 (i/N)·이어서 suffix.
    고정 캔버스(size!=None)에서만 동작 — auto(content-fit)는 캔버스가 늘어나므로 분할 불필요.
    """
    if size is None:
        return scenarios
    W, H = size
    avail_w = W - 2 * FIT_MARGIN
    avail_h = H - HEADER_BAND - FOOTER_BAND
    budget = avail_h / SEQ_PAGE_MIN_SCALE
    out = []
    for sc in scenarios:
        steps = sc.get("steps") or []
        full = R.layout_sequence(data, sc)
        s = min(avail_w / max(full["width"], 1), avail_h / max(full["height"], 1))
        if s >= SEQ_PAGE_MIN_SCALE or len(steps) <= 1:
            out.append(sc)
            continue
        chunks, cur = [], []
        for st in steps:
            cur.append(st)
            if R.layout_sequence(data, {**sc, "steps": cur})["height"] > budget and len(cur) > 1:
                chunks.append(cur[:-1])
                cur = [st]
        if cur:
            chunks.append(cur)
        base = {k: v for k, v in sc.items() if k != "steps"}
        n_pages = len(chunks)
        for i, ch in enumerate(chunks, 1):
            suffix = f" ({i}/{n_pages})" + (" · 이어서" if i > 1 else "") if n_pages > 1 else ""
            out.append({**base, "title": base.get("title", "") + suffix, "steps": ch})
    return out


def export_sequence(data, out_path, slide_size="wide", paginate=True):
    """sequence 뷰 JSON → .pptx (시나리오 1개 = 슬라이드 1장, 길면 자동 페이지 분할).

    render.py 의 layout_sequence() 기하를 그대로 소비 — actor 박스·라이프라인·
    activation bar·message(kind별 색)·note 를 px→EMU(×9525)로 배치(재구현 금지).
    self·note 는 render 와 동일 좌표에 배치하되 self-loop 는 엘보 커넥터로 근사.
    paginate=True(기본): 한 장에 넣으면 판독 불가하게 축소될 긴 시나리오를 여러 슬라이드로 나눔.
    """
    from pptx import Presentation
    from pptx.util import Emu, Pt
    from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.dml import MSO_LINE_DASH_STYLE
    from pptx.dml.color import RGBColor
    from pptx.oxml.ns import qn

    R = _load_render()
    size = _parse_slide_size(slide_size)
    scenarios = data.get("scenarios") or []
    if paginate:
        scenarios = _paginate_sequences(R, data, scenarios, size)
    layouts = [R.layout_sequence(data, sc) for sc in scenarios]
    max_w = max((L["width"] for L in layouts), default=400)
    max_h = max((L["height"] for L in layouts), default=300)
    SW, SH, s, dx, dy = _fit(max_w, max_h, size)
    emu = lambda px: Emu(int(round(px * s * PX_TO_EMU)))   # 크기·길이
    X = lambda px: emu(px + dx)                            # 위치(x) — 중앙 배치 오프셋
    Y = lambda px: emu(px + dy)                            # 위치(y) — 헤더 밴드 아래
    BOX_W, BOX_H, ACT_W, ZONE_H = R.BOX_W, R.BOX_H, R.ACT_W, R.ZONE_H

    prs = Presentation()
    prs.slide_width = Emu(SW * PX_TO_EMU)
    prs.slide_height = Emu(SH * PX_TO_EMU)
    blank = prs.slide_layouts[6]

    def _labels(shapes, x, y, w, h, specs, align, anchor):
        """specs: [(text, size_pt, rgb, bold)] → 문단별 텍스트박스."""
        tb = shapes.add_textbox(X(x), Y(y), emu(w), emu(h))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        for i, (text, size, rgb, bold) in enumerate(specs):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            r = p.add_run()
            r.text = text
            r.font.size = Pt(_fpt(size, s))
            r.font.color.rgb = RGBColor(*rgb)
            r.font.bold = bold

    def _arrow(conn, rgb):
        conn.line.color.rgb = RGBColor(*rgb)
        ln = conn.line._get_or_add_ln()
        ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle"}))

    for idx, (sc, L) in enumerate(zip(scenarios, layouts), 1):
        slide = prs.slides.add_slide(blank)
        shapes = slide.shapes

        _draw_chrome(shapes, SW, SH, data.get("system", ""), sc.get("title", ""),
                     data.get("source", ""), idx, len(layouts))

        zone_y, box_y, bottom = L["zone_y"], L["box_y"], L["bottom"]

        # 존 밴드(배경)
        for z in L["zones"]:
            zx1, zx2 = z["x1"], z["x2"]
            zb = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, X(zx1), Y(zone_y),
                                  emu(zx2 - zx1), emu(ZONE_H))
            zb.fill.solid()
            zb.fill.fore_color.rgb = RGBColor(*SEQ_ACTOR_FILL)
            zb.line.color.rgb = RGBColor(*SEQ_ACTOR_LINE)
            zb.shadow.inherit = False
            tf = zb.text_frame
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.text = z["name"]
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            tf.paragraphs[0].runs[0].font.size = Pt(_fpt(9, s))
            tf.paragraphs[0].runs[0].font.bold = True
            tf.paragraphs[0].runs[0].font.color.rgb = RGBColor(*SEQ_ACTOR_LINE)

        # 라이프라인(세로 점선)
        for a in L["actors"]:
            x = a["x"]
            ll = shapes.add_connector(MSO_CONNECTOR.STRAIGHT, X(x), Y(box_y + BOX_H),
                                      X(x), Y(bottom))
            ll.line.color.rgb = RGBColor(*SEQ_LIFELINE)
            ll.line.width = Pt(max(0.75, s))
            ll.line.dash_style = MSO_LINE_DASH_STYLE.DASH

        # 액티베이션 바
        for b in L["bars"]:
            ab = shapes.add_shape(MSO_SHAPE.RECTANGLE, X(b["x"]), Y(b["y"]),
                                  emu(ACT_W), emu(b["h"]))
            ab.fill.solid()
            ab.fill.fore_color.rgb = RGBColor(*SEQ_ACT_FILL)
            ab.line.color.rgb = RGBColor(*SEQ_ACT_LINE)
            ab.shadow.inherit = False

        # 액터 박스(포트/라인 2단)
        for a in L["actors"]:
            x = a["x"]
            box = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, X(x - BOX_W / 2), Y(box_y),
                                   emu(BOX_W), emu(BOX_H))
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(*SEQ_ACTOR_FILL)
            box.line.color.rgb = RGBColor(*SEQ_ACTOR_LINE)
            box.shadow.inherit = False
            tf = box.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            _fill_lines(tf, a["name"], s, size=11, sub_size=8,
                        color=SEQ_ACTOR_LINE, sub_color=SEQ_MUTED)
            if a["attrs"]:
                pp = tf.add_paragraph()
                pp.alignment = PP_ALIGN.CENTER
                r = pp.add_run()
                r.text = a["attrs"]
                r.font.size = Pt(_fpt(8, s))
                r.font.color.rgb = RGBColor(*SEQ_MUTED)

        # 메시지 / 노트
        for st in L["steps"]:
            if st["type"] == "note":
                x1, x2, y, h, lines = st["x1"], st["x2"], st["y"], st["h"], st["lines"]
                nb = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, X(x1), Y(y),
                                      emu(x2 - x1), emu(h))
                nb.fill.solid()
                nb.fill.fore_color.rgb = RGBColor(*SEQ_NOTE_FILL)
                nb.line.color.rgb = RGBColor(*SEQ_NOTE_LINE)
                nb.shadow.inherit = False
                if lines:
                    _labels(shapes, x1, y, x2 - x1, h,
                            [(ln, 9, SEQ_TEXT, False) for ln in lines],
                            PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
                continue

            kind, y, lines, mid = st["kind"], st["y"], st["lines"], st["mid"]
            rgb = SEQ_ARROW.get(kind, SEQ_ARROW["req"])
            bold = kind in ("req", "self")
            if st["self"]:
                sx = st["self_x"]
                conn = shapes.add_connector(MSO_CONNECTOR.ELBOW, X(sx), Y(y - 14),
                                            X(sx + 44), Y(y))
                _arrow(conn, rgb)
                lbl_x, lbl_align = sx + 52, PP_ALIGN.LEFT
            else:
                x1, x2 = st["x1"], st["x2"]
                conn = shapes.add_connector(MSO_CONNECTOR.STRAIGHT, X(x1), Y(y),
                                            X(x2), Y(y))
                _arrow(conn, rgb)
                lbl_x, lbl_align = mid - 70, PP_ALIGN.CENTER

            # 메시지 라벨(번호 인라인 포함) — 화살표 위
            if lines:
                lh = 15 * len(lines)
                _labels(shapes, lbl_x, y - 6 - lh, 140, lh + 4,
                        [(ln, 9, rgb, bold) for ln in lines],
                        lbl_align, MSO_ANCHOR.BOTTOM)
            # protocol / sub — 화살표 아래
            extra_specs = []
            for cls, val in st["extras"]:
                if cls == "proto":
                    extra_specs.append((f"( {val} )", 8, SEQ_MUTED, False))
                else:
                    extra_specs.append((val, 8, SEQ_WARN, False))
            if extra_specs:
                _labels(shapes, mid - 70, y + 2, 140, 15 * len(extra_specs) + 4,
                        extra_specs, PP_ALIGN.CENTER, MSO_ANCHOR.TOP)

    prs.save(str(out_path))
    return len(layouts)


# (exporter, render.py 검증기명) — 검증기는 render.py 단일 진실을 재사용한다.
_DISPATCH = {
    "sequence": (export_sequence, "validate"),
    "component": (export_component, "validate_component"),
    "topology": (export_topology, "validate_topology"),
}


def main():
    ap = argparse.ArgumentParser(description="flowcast 흐름도 → 편집가능 .pptx export (sequence·component·topology)")
    ap.add_argument("data", help="흐름도 뷰 JSON 경로 (view: sequence|component|topology)")
    ap.add_argument("-o", "--out", help="출력 .pptx (기본: 입력과 같은 위치 .pptx)")
    ap.add_argument("--slide-size", default="wide",
                    help="슬라이드 캔버스: wide(1920x1080, 기본)|auto(content-fit)|{W}x{H} px")
    ap.add_argument("--no-paginate", action="store_true",
                    help="긴 sequence 자동 페이지 분할 비활성화 (기본: 활성)")
    args = ap.parse_args()

    if not _import_pptx():
        print("python-pptx 가 필요합니다 (flowcast PPT export 전용 선택적 의존성).\n"
              "  pip install python-pptx", file=sys.stderr)
        return 2

    path = Path(args.data)
    if not path.exists():
        print(f"파일 없음: {path}", file=sys.stderr)
        return 1
    data = json.loads(path.read_text(encoding="utf-8"))
    # sequence 는 view 미지정이 기본값 → render.py 와 동일하게 sequence 로 간주
    view = data.get("view", "sequence")
    if view not in _DISPATCH:
        print(f"이 export 는 sequence·component·topology 뷰를 지원합니다 (view={view!r}).",
              file=sys.stderr)
        return 1

    # 렌더와 같은 검증을 먼저 통과시킨다 — 없으면 미정의 참조가 KeyError 트레이스백으로 샌다.
    exporter, validator_name = _DISPATCH[view]
    errors, warnings = getattr(_load_render(), validator_name)(data)
    for w in warnings:
        print(f"경고: {w}", file=sys.stderr)
    if errors:
        for e in errors:
            print(f"검증 오류: {e}", file=sys.stderr)
        return 1

    out = Path(args.out) if args.out else path.with_suffix(".pptx")
    kwargs = {"slide_size": args.slide_size}
    if view == "sequence":
        kwargs["paginate"] = not args.no_paginate
    n = exporter(data, out, **kwargs)
    print(f"pptx: {out} (슬라이드 {n})")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
