"""flowcast scripts/pptx_export.py — component 뷰 → .pptx round-trip 테스트.

python-pptx 는 export 전용 선택적 의존성 → 미설치 환경에선 이 파일 전체 skip.
검증 전략: export 한 .pptx 를 python-pptx 로 **다시 열어** 도형·텍스트·위치가 보존됐는지 확인
(= PowerPoint 편집가능성의 자동 프록시).
"""

import importlib.util
import json
from pathlib import Path

import pytest

pytest.importorskip("pptx")
from pptx import Presentation  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE_TYPE  # noqa: E402

_spec = importlib.util.spec_from_file_location(
    "pptx_export", Path(__file__).parent.parent / "scripts" / "pptx_export.py")
_mod = importlib.util.module_from_spec(_spec)
_spec.loader.exec_module(_mod)
export_component = _mod.export_component
export_topology = _mod.export_topology
export_sequence = _mod.export_sequence

EX = Path(__file__).parent.parent / "examples"
SRC = EX / "microservice-component.json"
TOPO = EX / "three-tier-topology.json"
SEQ = EX / "order-service-sequence.json"


def _export(tmp_path):
    data = json.loads(SRC.read_text(encoding="utf-8"))
    out = tmp_path / "c.pptx"
    n = export_component(data, out)
    return data, out, n


def _all_text(prs):
    parts = []
    for s in prs.slides:
        for sh in s.shapes:
            if sh.has_text_frame:
                parts.append(sh.text_frame.text)
            elif sh.has_table:   # 흐름 설명 네이티브 표 — 셀 텍스트도 포함
                for row in sh.table.rows:
                    for cell in row.cells:
                        parts.append(cell.text)
    return "\n".join(parts)


def test_slide_per_scenario(tmp_path):
    data, out, n = _export(tmp_path)
    assert n == len(data["scenarios"]) == 2
    assert len(Presentation(str(out)).slides) == 2


def test_node_names_and_ports_preserved(tmp_path):
    data, out, _ = _export(tmp_path)
    text = _all_text(Presentation(str(out)))
    for sc in data["scenarios"]:
        for nd in sc["nodes"]:
            assert nd["name"] in text
            if nd.get("port"):
                assert f"Port: {nd['port']}" in text


def test_edge_labels_preserved(tmp_path):
    _, out, _ = _export(tmp_path)
    text = _all_text(Presentation(str(out)))
    assert "(1)" in text and "주문 생성" in text and "( http )" in text
    assert "카드 승인" in text and "( https )" in text


def test_connector_count_equals_edges(tmp_path):
    data, out, _ = _export(tmp_path)
    prs = Presentation(str(out))
    total_edges = sum(len(sc.get("edges") or []) for sc in data["scenarios"])
    conns = sum(1 for s in prs.slides for sh in s.shapes
                if sh.shape_type == MSO_SHAPE_TYPE.LINE)
    assert conns == total_edges


def test_node_boxes_within_slide_bounds(tmp_path):
    data, out, _ = _export(tmp_path)
    prs = Presentation(str(out))
    names = {nd["name"] for sc in data["scenarios"] for nd in sc["nodes"]}
    total_nodes = sum(len(sc["nodes"]) for sc in data["scenarios"])
    W, H = prs.slide_width, prs.slide_height
    found = 0
    for s in prs.slides:
        for sh in s.shapes:
            if sh.has_text_frame and sh.text_frame.text.split("\n")[0] in names:
                assert 0 <= sh.left and sh.left + sh.width <= W
                assert 0 <= sh.top and sh.top + sh.height <= H
                found += 1
    assert found == total_nodes  # 모든 노드 박스가 슬라이드 경계 내


def test_non_component_view_returns_none_of_our_shapes(tmp_path):
    # export_component 는 component 전용 — view 가드는 main()에서. 여기선 정상 입력만 다룸.
    data, out, _ = _export(tmp_path)
    prs = Presentation(str(out))
    # 존 박스도 하나 존재해야(첫 시나리오 < Internal >)
    assert "< Internal >" in _all_text(prs)


# ── topology export ───────────────────────────────────────────

def _export_topo(tmp_path):
    data = json.loads(TOPO.read_text(encoding="utf-8"))
    out = tmp_path / "t.pptx"
    n = export_topology(data, out)
    return data, out, n


def test_topology_slide_per_scenario(tmp_path):
    data, out, n = _export_topo(tmp_path)
    assert n == len(data["scenarios"])
    assert len(Presentation(str(out)).slides) == n


def test_topology_nodes_and_zones_preserved(tmp_path):
    data, out, _ = _export_topo(tmp_path)
    text = _all_text(Presentation(str(out)))
    for nd in data["nodes"]:
        assert nd["name"] in text
    for z in data.get("zones") or []:
        assert z["name"] in text


def test_topology_segment_labels_preserved(tmp_path):
    _, out, _ = _export_topo(tmp_path)
    text = _all_text(Presentation(str(out)))
    assert "HTTPS 요청" in text and "부하 분산" in text  # slide2 segments


def test_topology_connector_count(tmp_path):
    data, out, _ = _export_topo(tmp_path)
    prs = Presentation(str(out))
    links = len(data.get("links") or [])
    ids = {n["id"] for n in data["nodes"]}
    seg_arrows = sum(
        1 for sc in data["scenarios"] for sg in (sc.get("segments") or [])
        if not sg.get("self") and sg.get("to") in ids)
    expected = links * len(data["scenarios"]) + seg_arrows  # 링크는 슬라이드마다 공통
    conns = sum(1 for s in prs.slides for sh in s.shapes
                if sh.shape_type == MSO_SHAPE_TYPE.LINE)
    assert conns == expected


# ── sequence export ───────────────────────────────────────────

def _export_seq(tmp_path):
    data = json.loads(SEQ.read_text(encoding="utf-8"))
    out = tmp_path / "s.pptx"
    n = export_sequence(data, out)
    return data, out, n


def test_sequence_slide_per_scenario(tmp_path):
    data, out, n = _export_seq(tmp_path)
    assert n == len(data["scenarios"]) == 2
    assert len(Presentation(str(out)).slides) == 2


def test_sequence_actors_and_zones_preserved(tmp_path):
    data, out, _ = _export_seq(tmp_path)
    text = _all_text(Presentation(str(out)))
    for a in data["actors"]:
        assert a["name"] in text
        if a.get("port"):
            assert str(a["port"]) in text   # sequence 서브라벨 = raw 값(render.py 동일)
    for z in data["zones"]:
        assert z["name"] in text


def test_sequence_message_labels_and_numbers(tmp_path):
    _, out, _ = _export_seq(tmp_path)
    text = _all_text(Presentation(str(out)))
    assert "1. 상품 주문" in text and "8. 주문 완료" in text  # 번호 인라인
    assert "( HTTPS )" in text                              # protocol extra


def test_sequence_connector_count(tmp_path):
    # 라이프라인(actor 수) + 메시지 커넥터(note 제외), 슬라이드별 합산
    data, out, _ = _export_seq(tmp_path)
    prs = Presentation(str(out))
    n_actors = len(data["actors"])
    expected = sum(n_actors + sum(1 for st in sc["steps"] if st["kind"] != "note")
                   for sc in data["scenarios"])
    conns = sum(1 for s in prs.slides for sh in s.shapes
                if sh.shape_type == MSO_SHAPE_TYPE.LINE)
    assert conns == expected


def test_sequence_actor_boxes_within_bounds(tmp_path):
    data, out, _ = _export_seq(tmp_path)
    prs = Presentation(str(out))
    names = {a["name"] for a in data["actors"]}
    W, H = prs.slide_width, prs.slide_height
    found = 0
    for s in prs.slides:
        for sh in s.shapes:
            if sh.has_text_frame and sh.text_frame.text.split("\n")[0] in names:
                assert 0 <= sh.left and sh.left + sh.width <= W
                assert 0 <= sh.top and sh.top + sh.height <= H
                found += 1
    assert found == len(names) * len(data["scenarios"])


def test_sequence_self_and_note_render(tmp_path):
    # self·note 는 예제에 없어 합성 데이터로 두 분기 커버
    data = {
        "system": "S",
        "actors": [{"id": "a", "name": "Alpha"}, {"id": "b", "name": "Beta"}],
        "scenarios": [{"title": "T", "steps": [
            {"n": 1, "from": "a", "to": "b", "label": "요청", "kind": "req"},
            {"from": "b", "to": "b", "label": "내부 처리", "kind": "self"},
            {"from": "a", "to": "b", "label": "메모", "kind": "note"},
        ]}],
    }
    out = tmp_path / "sn.pptx"
    assert export_sequence(data, out) == 1
    text = _all_text(Presentation(str(out)))
    assert "1. 요청" in text and "내부 처리" in text and "메모" in text


# ── #17: 멀티라인 스타일 · topology 배지/범례 · 슬라이드 캔버스 ──

EMU_PX = 9525


def test_default_slide_size_wide_16_9(tmp_path):
    for fn, src in ((export_component, SRC), (export_topology, TOPO), (export_sequence, SEQ)):
        data = json.loads(src.read_text(encoding="utf-8"))
        out = tmp_path / f"w-{src.stem}.pptx"
        fn(data, out)
        prs = Presentation(str(out))
        assert prs.slide_width == 1920 * EMU_PX
        assert prs.slide_height == 1080 * EMU_PX


def test_slide_size_auto_is_content_fit(tmp_path):
    data = json.loads(TOPO.read_text(encoding="utf-8"))
    out = tmp_path / "a.pptx"
    export_topology(data, out, slide_size="auto")
    prs = Presentation(str(out))
    assert (prs.slide_width, prs.slide_height) != (1920 * EMU_PX, 1080 * EMU_PX)


# ── 상단 타이틀 밴드 (render.py hero 헤더 대응: eyebrow + h1(scenario title) + sub(source)) ──

def test_title_header_component(tmp_path):
    data, out, _ = _export(tmp_path)
    text = _all_text(Presentation(str(out)))
    assert f"{data['system']} · IF 흐름도" in text
    for sc in data["scenarios"]:
        assert sc["title"] in text
    if data.get("source"):
        assert data["source"] in text


def test_title_header_topology(tmp_path):
    data, out, _ = _export_topo(tmp_path)
    text = _all_text(Presentation(str(out)))
    assert f"{data['system']} · IF 흐름도" in text
    for sc in data["scenarios"]:
        assert sc["title"] in text


def test_title_header_sequence(tmp_path):
    data, out, _ = _export_seq(tmp_path)
    text = _all_text(Presentation(str(out)))
    assert f"{data['system']} · IF 흐름도" in text
    for sc in data["scenarios"]:
        assert sc["title"] in text


def test_title_header_auto_mode(tmp_path):
    data = json.loads(SRC.read_text(encoding="utf-8"))
    out = tmp_path / "ca.pptx"
    export_component(data, out, slide_size="auto")
    text = _all_text(Presentation(str(out)))
    assert f"{data['system']} · IF 흐름도" in text


def test_title_band_above_content(tmp_path):
    # 타이틀 밴드는 콘텐츠(노드) 위에 예약된다 — 겹침 방지
    data, out, _ = _export(tmp_path)
    slide = Presentation(str(out)).slides[0]
    node_names = {nd["name"].split("\n")[0] for sc in data["scenarios"] for nd in sc["nodes"]}
    title_tops = [sh.top for sh in slide.shapes
                  if sh.has_text_frame and "IF 흐름도" in sh.text_frame.text]
    node_tops = [sh.top for sh in slide.shapes
                 if sh.has_text_frame and sh.text_frame.text.split("\n")[0] in node_names]
    assert title_tops and node_tops
    assert min(title_tops) < min(node_tops)


def test_slide_size_custom(tmp_path):
    data = json.loads(SEQ.read_text(encoding="utf-8"))
    out = tmp_path / "c1280.pptx"
    export_sequence(data, out, slide_size="1280x720")
    prs = Presentation(str(out))
    assert prs.slide_width == 1280 * EMU_PX and prs.slide_height == 720 * EMU_PX


def test_multiline_node_name_all_runs_styled(tmp_path):
    # 멀티라인 이름 둘째 줄이 템플릿 기본값(18pt)으로 새지 않아야 — 전 문단 명시 스타일
    data = json.loads(TOPO.read_text(encoding="utf-8"))
    data["nodes"][0]["name"] = "Alpha\n(부가 설명)"
    out = tmp_path / "m.pptx"
    export_topology(data, out)
    prs = Presentation(str(out))
    target = [sh for s in prs.slides for sh in s.shapes
              if sh.has_text_frame and sh.text_frame.text.startswith("Alpha")]
    assert target
    for sh in target:
        paras = sh.text_frame.paragraphs
        assert len(paras) == 2
        for p in paras:
            assert p.runs and all(r.font.size is not None for r in p.runs)
        assert paras[1].runs[0].font.size < paras[0].runs[0].font.size


def test_topology_badges_and_legend(tmp_path):
    # 세그먼트 라벨 전문은 엣지 중점이 아니라 번호 배지 + 하단 범례로
    data, out, _ = _export_topo(tmp_path)
    prs = Presentation(str(out))
    text = _all_text(prs)
    assert "흐름 설명" in text
    for i, sc in enumerate(data["scenarios"]):
        numbered = [sg for sg in (sc.get("segments") or []) if sg.get("n") is not None]
        badges = [sh for sh in prs.slides[i].shapes
                  if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and sh.has_text_frame
                  and sh.text_frame.text.strip().isdigit() and sh.name != "legbadge"]
        assert len(badges) == len(numbered)
        for sg in numbered:
            if sg.get("label"):
                # 범례는 [배지 | 단계 | 설명] 표 — label 의 단계/설명 텍스트가 본문에 존재
                assert sg["label"][:15] in text


def test_topology_legend_native_table(tmp_path):
    # 흐름 설명 범례 = 네이티브 표(GraphicFrame) — 헤더 [# | 단계 | 설명 · 기술] + 세그먼트당 1행
    data, out, _ = _export_topo(tmp_path)
    prs = Presentation(str(out))
    for i, sc in enumerate(data["scenarios"]):
        labelled = [sg for sg in (sc.get("segments") or []) if sg.get("label")]
        tables = [sh.table for sh in prs.slides[i].shapes if sh.has_table]
        if not labelled:
            assert not tables            # 라벨 없으면 범례 표 없음
            continue
        assert tables
        for t in tables:
            assert [t.cell(0, c).text for c in range(3)] == ["#", "단계", "설명 · 기술"]
        data_rows = sum(len(t.rows) - 1 for t in tables)   # 헤더 제외
        assert data_rows == len(labelled)
        cell_text = "\n".join(t.cell(r, c).text
                              for t in tables for r in range(len(t.rows)) for c in range(3))
        for sg in labelled:
            assert sg["label"][:12] in cell_text


def test_font_scale_follows_canvas(tmp_path):
    # 작은 예제 → wide 캔버스 업스케일 — 노드 폰트가 base(10pt)보다 커져야
    data = json.loads(TOPO.read_text(encoding="utf-8"))
    out = tmp_path / "f.pptx"
    export_topology(data, out)
    prs = Presentation(str(out))
    node_names = {n["name"] for n in data["nodes"]}
    sizes = [sh.text_frame.paragraphs[0].runs[0].font.size.pt
             for s in prs.slides for sh in s.shapes
             if sh.has_text_frame and sh.text_frame.text.split("\n")[0] in node_names]
    assert sizes and all(sz > 10 for sz in sizes)


def test_topology_badge_overlap_spread(tmp_path):
    # 동일 엣지 2개 → 배지 완전 중첩 — spread 로 지름(22px) 이상 분리 (#19). auto=scale 1
    data = {
        "view": "topology",
        "system": "S",
        "nodes": [
            {"id": "a", "name": "A", "col": 0, "row": 0},
            {"id": "b", "name": "B", "col": 1, "row": 0},
        ],
        "scenarios": [{"title": "T", "segments": [
            {"n": 1, "from": "a", "to": "b", "label": "one"},
            {"n": 2, "from": "a", "to": "b", "label": "two"},
        ]}],
    }
    out = tmp_path / "spread.pptx"
    export_topology(data, out, slide_size="auto")
    prs = Presentation(str(out))
    centers = [(sh.left + sh.width / 2, sh.top + sh.height / 2)
               for sh in prs.slides[0].shapes
               if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and sh.has_text_frame
               and sh.text_frame.text.strip().isdigit() and sh.name != "legbadge"]
    assert len(centers) == 2
    d = ((centers[0][0] - centers[1][0]) ** 2 + (centers[0][1] - centers[1][1]) ** 2) ** 0.5
    assert d >= 22 * 9525   # 지름 22px 이상 (EMU)


# ── #25: z-순서 패리티 — 노드가 커넥터를 가림 (HTML과 동일) ──

def _zorder_conn_under_nodes(prs, node_first_lines):
    """모든 슬라이드에서 LINE(커넥터) 도형이 노드 박스보다 먼저(아래) 삽입됐는지."""
    for slide in prs.slides:
        idx_conn = [i for i, sh in enumerate(slide.shapes)
                    if sh.shape_type == MSO_SHAPE_TYPE.LINE]
        idx_node = [i for i, sh in enumerate(slide.shapes)
                    if sh.has_text_frame and sh.text_frame.text.split("\n")[0] in node_first_lines]
        if idx_conn and idx_node:
            assert max(idx_conn) < min(idx_node)


def test_topology_zorder_nodes_over_connectors(tmp_path):
    data, out, _ = _export_topo(tmp_path)
    prs = Presentation(str(out))
    names = {str(n["name"]).split("\n")[0] for n in data["nodes"]}
    _zorder_conn_under_nodes(prs, names)


def test_component_zorder_nodes_over_connectors(tmp_path):
    data, out, _ = _export(tmp_path)
    prs = Presentation(str(out))
    names = {str(nd["name"]).split("\n")[0] for sc in data["scenarios"] for nd in sc["nodes"]}
    _zorder_conn_under_nodes(prs, names)


# ── #41: 상단 제목영역 일관성(뷰 무관 고정 앵커) + 제목/본문·본문/푸터 구분선 + 페이지 ──

_VIEWS = ((_mod.export_component, SRC), (_mod.export_topology, TOPO),
          (_mod.export_sequence, SEQ))


def _first_slide(fn, src, tmp_path):
    data = json.loads(src.read_text(encoding="utf-8"))
    out = tmp_path / f"h-{src.stem}.pptx"
    fn(data, out)
    return Presentation(str(out)).slides[0]


def _title_box(slide):
    return next(sh for sh in slide.shapes
               if sh.has_text_frame and "IF 흐름도" in sh.text_frame.text)


def _rules(slide):
    """제목/본문·본문/푸터 구분선 = 얇은(≤3px) 텍스트 없는 사각형(AUTO_SHAPE)."""
    return [sh for sh in slide.shapes
            if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
            and not sh.text_frame.text.strip() and sh.height <= 3 * EMU_PX]


def test_header_anchor_consistent_across_views(tmp_path):
    # 제목 좌상단이 sequence·component·topology 에서 동일 (fit-centering 아티팩트 제거)
    boxes = [_title_box(_first_slide(fn, src, tmp_path)) for fn, src in _VIEWS]
    assert len({b.left for b in boxes}) == 1
    assert len({b.top for b in boxes}) == 1


def test_header_h1_font_consistent_across_views(tmp_path):
    # h1 = eyebrow 다음 문단. 콘텐츠 배율 s 와 무관하게 뷰 간 동일해야
    sizes = {_title_box(_first_slide(fn, src, tmp_path))
             .text_frame.paragraphs[1].runs[0].font.size.pt for fn, src in _VIEWS}
    assert len(sizes) == 1


def test_divider_rules_are_shapes_not_connectors(tmp_path):
    # 구분선 2개(제목/본문·본문/푸터). LINE(커넥터) 아니어야 커넥터 수 테스트 불변
    for fn, src in _VIEWS:
        rules = _rules(_first_slide(fn, src, tmp_path))
        assert len(rules) == 2
        assert all(r.shape_type != MSO_SHAPE_TYPE.LINE for r in rules)


def test_footer_page_number_present(tmp_path):
    data = json.loads(TOPO.read_text(encoding="utf-8"))
    out = tmp_path / "pg.pptx"
    n = export_topology(data, out)
    assert f"1 / {n}" in _all_text(Presentation(str(out)))


def test_content_below_header_band(tmp_path):
    # 모든 노드가 헤더 밴드(구분선) 아래 — 제목/본문 겹침 없음
    slide = _first_slide(export_component, SRC, tmp_path)
    rule_y = min(r.top for r in _rules(slide))
    data = json.loads(SRC.read_text(encoding="utf-8"))
    names = {nd["name"].split("\n")[0] for sc in data["scenarios"] for nd in sc["nodes"]}
    node_tops = [sh.top for sh in slide.shapes
                 if sh.has_text_frame and sh.text_frame.text.split("\n")[0] in names]
    assert node_tops and min(node_tops) >= rule_y


# ── #43: 긴 sequence 자동 페이지네이션 (읽히는 범위 + 다음 슬라이드 연장) ──

def _long_seq(n):
    return {
        "system": "S",
        "actors": [{"id": "a", "name": "Alpha"}, {"id": "b", "name": "Beta"},
                   {"id": "c", "name": "Gamma"}],
        "scenarios": [{"title": "Long", "steps": [
            {"n": i, "from": ("a" if i % 2 else "b"), "to": ("b" if i % 2 else "c"),
             "label": f"step {i}", "kind": "req", "protocol": "HTTPS"}
            for i in range(1, n + 1)]}],
    }


def test_long_sequence_paginates(tmp_path):
    out = tmp_path / "lng.pptx"
    n = export_sequence(_long_seq(30), out)
    assert n > 1
    assert len(Presentation(str(out)).slides) == n


def test_short_sequence_single_slide(tmp_path):
    assert export_sequence(_long_seq(4), tmp_path / "sht.pptx") == 1


def test_pagination_can_be_disabled(tmp_path):
    # paginate=False → 길어도 한 장(기존 축소 동작)
    assert export_sequence(_long_seq(30), tmp_path / "nop.pptx", paginate=False) == 1


def test_pagination_skipped_in_auto_mode(tmp_path):
    # auto=content-fit → 캔버스가 늘어나 분할 불필요
    assert export_sequence(_long_seq(30), tmp_path / "au.pptx", slide_size="auto") == 1


def test_paginated_actors_repeat_each_slide(tmp_path):
    # 액터 헤더·라이프라인이 매 슬라이드 재렌더(연장선)
    out = tmp_path / "rep.pptx"
    export_sequence(_long_seq(30), out)
    prs = Presentation(str(out))
    for slide in prs.slides:
        text = "\n".join(sh.text_frame.text for sh in slide.shapes if sh.has_text_frame)
        assert all(name in text for name in ("Alpha", "Beta", "Gamma"))


def test_paginated_step_numbers_continuous(tmp_path):
    # step.n 유지 → 페이지 넘어 번호 연속 (첫·끝 모두 존재)
    out = tmp_path / "num.pptx"
    export_sequence(_long_seq(30), out)
    text = _all_text(Presentation(str(out)))
    assert "1. step 1" in text and "30. step 30" in text


def test_paginated_title_has_page_index(tmp_path):
    out = tmp_path / "idx.pptx"
    n = export_sequence(_long_seq(30), out)
    text = _all_text(Presentation(str(out)))
    assert f"(1/{n})" in text and "이어서" in text
